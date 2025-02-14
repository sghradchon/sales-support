import json
import os

import boto3
from pptx import Presentation
from pptx.util import Inches
from pathlib import Path
from pptx.enum.text import PP_ALIGN

S3_BUCKET = os.getenv("S3_BUCKET_NAME")
s3_client = boto3.client('s3')
TMP_DIR = "/tmp"


def handler(event, context):
    body = json.loads(event['body'])
    print("body:",body)
    paths = body.get('paths')
    print(" paths:",paths)
    
    response_body = {"slide_path":"hoge"}

    return {
        'statusCode': 200,
        'headers': {
            'Access-Control-Allow-Headers': '*',
            'Access-Control-Allow-Origin': '*',
            'Access-Control-Allow-Methods': 'OPTIONS,POST,GET'
        },
        'body': json.dumps(response_body)
    }
    

def download_from_s3(s3_path):
    """S3 からファイルをダウンロードし、一時フォルダに保存"""
    local_path = os.path.join(TMP_DIR, os.path.basename(s3_path))
    s3_client.download_file(S3_BUCKET, s3_path, local_path)
    return local_path

def get_parts_list(product_img_path):
    """製品画像のパスから対応する部品画像リストを取得（S3 を検索）"""
    product_name = Path(product_img_path).stem  # `productA.png` → `productA`
    parts_folder = f"parts-pictures/{product_name}/"

    response = s3_client.list_objects_v2(Bucket=S3_BUCKET, Prefix=parts_folder)
    part_images = []

    if "Contents" in response:
        for item in response["Contents"]:
            if item["Key"].lower().endswith((".jpg", ".jpeg", ".png", ".gif", ".bmp", ".tiff")):
                part_images.append(item["Key"])

    return product_name, part_images

def add_combined_slide(prs, product_name, product_img, part_imgs):
    """製品と部品を1枚のスライドに収めるように動的に配置"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # タイトルなしのスライド

    # ✅ タイトルテキストボックスを明示的に作成し、横幅を確保
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.2), Inches(8), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = product_name
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # 中央揃え
    title_frame.word_wrap = False  # 改行を防ぐ

    # ✅ 製品画像をダウンロード & 配置
    product_img_path = download_from_s3(product_img)
    prod_left, prod_top, prod_width, prod_height = Inches(1.0), Inches(1.0), Inches(7.5), Inches(3.5)
    slide.shapes.add_picture(product_img_path, prod_left, prod_top, prod_width, prod_height)

    # ✅ 部品画像の配置
    part_width, part_height = Inches(1.5), Inches(1.5)
    margin_x, margin_y = Inches(0.3), Inches(0.3)
    start_x, start_y = Inches(0.5), Inches(4.0)
    slide_width = Inches(10)
    available_height = Inches(7.5) - start_y  # 部品を配置できる最大の高さ

    max_cols = int((slide_width - start_x) / (part_width + margin_x))
    max_rows = len(part_imgs) // max_cols + (1 if len(part_imgs) % max_cols else 0)

    while (max_rows * (part_height + margin_y) > available_height) and prod_height > Inches(2.5):
        prod_height -= Inches(0.3)
        start_y += Inches(0.2)
        available_height = Inches(7.5) - start_y
        part_height = part_width = max(Inches(1.0), part_height - Inches(0.2))
        max_rows = len(part_imgs) // max_cols + (1 if len(part_imgs) % max_cols else 0)

    # ✅ 部品画像をダウンロード & 配置
    for i, part_img in enumerate(part_imgs):
        row, col = i // max_cols, i % max_cols
        left, top = start_x + (part_width + margin_x) * col, start_y + (part_height + margin_y) * row
        part_img_path = download_from_s3(part_img)
        slide.shapes.add_picture(part_img_path, left, top, part_width, part_height)

def create_pptx(product_image_paths):
    """受け取った製品画像リストを使って PowerPoint を作成"""
    prs = Presentation()

    for product_img_path in product_image_paths:
        product_name, part_images = get_parts_list(product_img_path)
        add_combined_slide(prs, product_name, product_img_path, part_images)

    pptx_filename = os.path.join(TMP_DIR, "output_slide.pptx")
    prs.save(pptx_filename)
    return pptx_filename

def upload_to_s3(file_path, s3_key):
    """S3 にファイルをアップロード"""
    s3_client.upload_file(file_path, S3_BUCKET, s3_key, ExtraArgs={"ContentType": "application/vnd.openxmlformats-officedocument.presentationml.presentation"})
    s3_url = f"https://{S3_BUCKET}.s3.amazonaws.com/{s3_key}"
    return s3_url

def handler(event, context):
    """Lambda のエントリポイント"""
    try:
        # ✅ リクエストボディの解析
        body = json.loads(event["body"])
        product_image_paths = body.get("product_images", [])

        if not product_image_paths:
            return {"statusCode": 400, "body": json.dumps({"error": "製品画像が指定されていません"})}

        # ✅ PowerPoint を作成
        pptx_path = create_pptx(product_image_paths)

        # ✅ S3 にアップロード
        s3_key = "slides/output_slide.pptx"
        s3_url = upload_to_s3(pptx_path, s3_key)
        response_body = {"message":"success","slide_url":s3_url}
        
        return {
            'statusCode': 200,
            'headers': {
                'Access-Control-Allow-Headers': '*',
                'Access-Control-Allow-Origin': '*',
                'Access-Control-Allow-Methods': 'OPTIONS,POST,GET'
            },
            'body': json.dumps(response_body)
        }
    

    except Exception as e:
        return {"statusCode": 500, "body": json.dumps({"error": str(e)})}